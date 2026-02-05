from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import uuid
from typing import List
from PIL import Image
import io
from .base import BaseExtractor, ContentItem
from ..config import IMAGE_STORE_DIR

class PPTXExtractor(BaseExtractor):
    def extract(self, file_path: str) -> List[ContentItem]:
        items = []
        prs = Presentation(file_path)
        filename = os.path.basename(file_path)
        
        for slide_index, slide in enumerate(prs.slides):
            page_num = slide_index + 1
            slide_text = []
            
            for shape in slide.shapes:
                # --- Text Extraction ---
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                
                # --- Image Extraction ---
                # Check for picture shape
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image_blob = shape.image.blob
                        pil_image = Image.open(io.BytesIO(image_blob))
                        
                        # Filter tiny images
                        if pil_image.width < 100 or pil_image.height < 100:
                            continue
                        
                        image_id = f"{uuid.uuid4()}.png"
                        save_path = os.path.join(IMAGE_STORE_DIR, image_id)
                        pil_image.save(save_path)
                        
                        items.append(ContentItem(
                            content=f"[[Image extracted from Slide {page_num}]]",
                            type="image",
                            source=filename,
                            page_num=page_num,
                            image_path=save_path,
                            metadata={
                                "width": pil_image.width, 
                                "height": pil_image.height,
                                "left": shape.left,
                                "top": shape.top
                            }
                        ))
                    except Exception as e:
                        print(f"Failed to extract image on slide {page_num}: {e}")
            
            # Consolidate text per slide
            if slide_text:
                full_slide_text = "\n".join(slide_text)
                items.append(ContentItem(
                    content=full_slide_text,
                    type="text",
                    source=filename,
                    page_num=page_num
                ))
                
        return items
